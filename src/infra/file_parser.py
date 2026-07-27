"""文件解析工具

策略（轻量级 dispatcher）：
- 每种格式独立解析器，全部使用纯 Python 库，零 ML 依赖
- PDF → pdfplumber
- DOCX → python-docx（段落 + 表格）
- PPTX → python-pptx（slide 文本）
- XLSX → openpyxl（按 sheet / 行 / 列遍历）
- XLS → xlrd（仅文本）
- CSV → 标准库 csv
- HTML / HTM → beautifulsoup4 get_text
- MD → markdown lib → bs4 get_text（保留结构）
- TXT → 编码探测
- 图片 → 占位（OCR 待云接入）

支持格式：PDF / DOCX / PPTX / XLSX / XLS / CSV / HTML / HTM / MD / TXT / PNG / JPG / JPEG
"""
import csv
import io
import logging
import re

logger = logging.getLogger(__name__)


_SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "pptx", "xlsx", "xls", "csv",
    "html", "htm", "md", "txt",
    "png", "jpg", "jpeg",
}


def parse_file(file_bytes: bytes, filename: str) -> str:
    """解析文件为文本。

    Args:
        file_bytes: 文件字节内容
        filename: 文件名（用于判断类型）

    Returns:
        提取的文本内容，解析失败或不支持返回空字符串
    """
    if not file_bytes:
        return ""

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in _SUPPORTED_EXTENSIONS:
        logger.warning(f"不支持的文件类型: {ext}, 文件名: {filename}")
        return ""

    parser = _PARSERS.get(ext)
    if parser is None:
        logger.warning(f"未实现解析器: {ext}")
        return ""

    try:
        result = parser(file_bytes)
        return result.strip() if result else ""
    except Exception as e:
        logger.error(f"解析失败 ({filename}): {e}")
        return ""


def get_supported_extensions() -> set[str]:
    return _SUPPORTED_EXTENSIONS.copy()


def is_supported(filename: str) -> bool:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    return ext in _SUPPORTED_EXTENSIONS


# ═══════════════════════════════════════════════════════════════════════════════
# 解析器实现
# ═══════════════════════════════════════════════════════════════════════════════


def _parse_pdf(file_bytes: bytes) -> str:
    """PDF 解析（pdfplumber）。"""
    import pdfplumber

    texts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n".join(texts)


def _parse_docx(file_bytes: bytes) -> str:
    """DOCX 解析：段落 + 表格。"""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for para in doc.paragraphs:
        if para.text and para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append("\n".join(rows))

    return "\n".join(parts)


def _parse_pptx(file_bytes: bytes) -> str:
    """PPTX 解析：按 slide 顺序提取所有 shape 文本。"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def _parse_xlsx(file_bytes: bytes) -> str:
    """XLSX 解析：按 sheet 输出。"""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _parse_xls(file_bytes: bytes) -> str:
    """XLS 解析（旧格式，仅文本）。"""
    import xlrd

    wb = xlrd.open_workbook(file_contents=file_bytes)
    parts = []
    for sheet in wb.sheets():
        rows = []
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col)).strip()
                     for col in range(sheet.ncols)]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _parse_csv(file_bytes: bytes) -> str:
    """CSV 解析（自动探测编码与分隔符）。"""
    text = _decode_text(file_bytes)
    if not text:
        return ""

    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect=dialect)
    rows = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
    return "\n".join(rows)


def _parse_html(file_bytes: bytes) -> str:
    """HTML 解析（bs4 get_text，保留段落结构）。"""
    from bs4 import BeautifulSoup

    text = _decode_text(file_bytes)
    if not text:
        return ""
    soup = BeautifulSoup(text, "lxml")
    for tag in soup(["script", "style", "meta", "link", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_md(file_bytes: bytes) -> str:
    """Markdown 解析：转 HTML → bs4 get_text，保留基本结构。"""
    from bs4 import BeautifulSoup
    import markdown as md_lib

    text = _decode_text(file_bytes)
    if not text:
        return ""
    html = md_lib.markdown(text, extensions=["extra"])
    soup = BeautifulSoup(html, "lxml")
    return soup.get_text(separator="\n", strip=True)


def _parse_txt(file_bytes: bytes) -> str:
    """纯文本解析（自动检测编码）。"""
    return _decode_text(file_bytes)


def _parse_image_ocr(file_bytes: bytes) -> str:
    """图片 OCR 占位（待云 OCR 接入）。

    Returns:
        空字符串 + warning 日志。
    """
    logger.warning("图片 OCR 暂未实现（待云 OCR 接入），文件大小: %d bytes", len(file_bytes))
    return ""


# ═══════════════════════════════════════════════════════════════════════════════
# 内部工具
# ═══════════════════════════════════════════════════════════════════════════════


_PARSERS = {
    "pdf":  _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "xlsx": _parse_xlsx,
    "xls":  _parse_xls,
    "csv":  _parse_csv,
    "html": _parse_html,
    "htm":  _parse_html,
    "md":   _parse_md,
    "txt":  _parse_txt,
    "png":  _parse_image_ocr,
    "jpg":  _parse_image_ocr,
    "jpeg": _parse_image_ocr,
}


def _decode_text(file_bytes: bytes) -> str:
    """编码探测：utf-8 → gbk → gb2312 → big5 → latin-1。"""
    for enc in ("utf-8", "gbk", "gb2312", "big5", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return file_bytes.decode("utf-8", errors="ignore")